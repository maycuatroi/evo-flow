import json
import pathlib

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation
from tqdm import tqdm

from evoflow.entities.data_manipulate.file_operator.file import File


def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            # pylint: disable=no-member
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape


class PPTXFile(File):
    def save(self, file_path=None) -> str:
        self.data.save(file_path)

    def get_info(self) -> str:
        self.data: Presentation

        properties = []
        for attribute in dir(self.data.core_properties):
            if not attribute.startswith("_"):
                properties.append(attribute)

        properties_dict = {"slides": len(self.data.slides)}
        for property_name in properties:
            value = getattr(self.data.core_properties, property_name)

            if (
                value is not None
                and value.__class__.__name__ in ["int", "str", "datetime"]
                and len(str(value)) > 0
            ):
                properties_dict[property_name] = str(value)

        return json.dumps(properties_dict, ensure_ascii=False, indent=2)

    def extract_images(self, output_path="data"):
        index = 0
        for picture in tqdm(iter_picture_shapes(self.data), desc="Extract Images"):
            image = picture.image
            # ---get image "file" contents---
            image_bytes = image.blob
            # ---make up a name for the file, e.g. 'image.jpg'---

            image_dir = f"{output_path}/{self.data.core_properties.title}"
            image_filename = f"{image_dir}/image_{index}.{image.ext}"

            pathlib.Path(image_dir).mkdir(parents=True, exist_ok=True)

            with open(image_filename, "wb") as file:
                file.write(image_bytes)
            index += 1

    def get_texts(self):
        prs = self.data
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

        return text_runs

    def __init__(self, **args):
        super().__init__(**args)
        self.data: Presentation
